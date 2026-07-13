import csv
import io
import json
import logging
from typing import Any

logger = logging.getLogger(__name__)


def render_csv(content: Any) -> bytes:
    buf = io.StringIO()
    rows = json.loads(content) if isinstance(content, str) and content.strip().startswith("[") else content
    if isinstance(rows, list) and rows and isinstance(rows[0], dict):
        writer = csv.DictWriter(buf, fieldnames=list(rows[0].keys()))
        writer.writeheader()
        writer.writerows(rows)
    elif isinstance(rows, str):
        buf.write(rows)
    else:
        buf.write(str(rows))
    return buf.getvalue().encode("utf-8")


def render_xlsx(content: Any) -> bytes:
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl not installed — run: pip install openpyxl")

    rows = json.loads(content) if isinstance(content, str) else content
    wb = openpyxl.Workbook()
    ws = wb.active
    if isinstance(rows, list) and rows and isinstance(rows[0], dict):
        headers = list(rows[0].keys())
        ws.append(headers)
        for row in rows:
            ws.append([row.get(h) for h in headers])
    elif isinstance(rows, list):
        for row in rows:
            ws.append(row if isinstance(row, list) else [row])

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def render_docx(content: str) -> bytes:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx not installed — run: pip install python-docx")

    doc = Document()
    for line in str(content).split("\n"):
        line = line.strip()
        if line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line:
            doc.add_paragraph(line)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def render_pptx(content: Any) -> bytes:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx not installed — run: pip install python-pptx")

    slides = json.loads(content) if isinstance(content, str) else content
    prs = Presentation()
    layout = prs.slide_layouts[1]

    for slide_data in (slides if isinstance(slides, list) else []):
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_data.get("content", "")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_pdf(content: str) -> bytes:
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
    except ImportError:
        raise RuntimeError("reportlab not installed — run: pip install reportlab")

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    for line in str(content).split("\n"):
        line = line.strip()
        if not line:
            story.append(Spacer(1, 10))
        elif line.startswith("### "):
            story.append(Paragraph(line[4:], styles["Heading3"]))
        elif line.startswith("## "):
            story.append(Paragraph(line[3:], styles["Heading2"]))
        elif line.startswith("# "):
            story.append(Paragraph(line[2:], styles["Heading1"]))
        else:
            story.append(Paragraph(line, styles["Normal"]))

    doc.build(story)
    return buf.getvalue()
